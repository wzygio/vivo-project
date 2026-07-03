from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(r"D:\wzy\Python\vivo-project\projects\special_projects_cpk_products")
EXPORT_DIR = PROJECT_DIR / "exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


COLORS = {
    "bg": RGBColor(249, 251, 254),
    "white": RGBColor(255, 255, 255),
    "blue": RGBColor(91, 155, 213),
    "deep_blue": RGBColor(68, 114, 196),
    "navy": RGBColor(61, 76, 96),
    "orange": RGBColor(237, 125, 49),
    "green": RGBColor(112, 173, 71),
    "yellow": RGBColor(255, 192, 0),
    "grey": RGBColor(166, 166, 166),
    "text": RGBColor(68, 84, 106),
    "light_line": RGBColor(214, 226, 241),
    "soft_blue": RGBColor(235, 244, 253),
    "soft_orange": RGBColor(255, 243, 228),
    "soft_green": RGBColor(237, 247, 235),
}


def set_font(run, size: int, color=COLORS["text"], bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def set_text(shape, text: str, size: int, color=COLORS["text"], bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)


def add_text(slide, text: str, x, y, w, h, size=12, color=COLORS["text"], bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, w, h)
    set_text(shape, text, size, color, bold, align)
    return shape


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.6)
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.8)
    return shape


def add_line(slide, x1, y1, x2, y2, color, width=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_header(slide, section: str, title: str, subtitle: str, page_no: str) -> None:
    add_text(slide, section, Inches(0.62), Inches(0.28), Inches(2.2), Inches(0.25), 9, COLORS["blue"], True)
    add_text(slide, title, Inches(0.62), Inches(0.55), Inches(8.0), Inches(0.42), 20, COLORS["text"], True)
    add_text(slide, subtitle, Inches(0.62), Inches(0.96), Inches(10.2), Inches(0.24), 8, COLORS["grey"])
    add_line(slide, Inches(0.62), Inches(1.25), Inches(12.6), Inches(1.25), COLORS["light_line"], 1)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), COLORS["blue"], COLORS["blue"])
    add_text(slide, page_no, Inches(12.2), Inches(6.95), Inches(0.6), Inches(0.24), 8, COLORS["grey"], False, PP_ALIGN.RIGHT)


def style_table(table, header_fill=COLORS["deep_blue"], font_size=6) -> None:
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_fill
                color = COLORS["white"]
                bold = True
                align = PP_ALIGN.CENTER
            else:
                cell.fill.fore_color.rgb = COLORS["white"] if row_idx % 2 else COLORS["bg"]
                color = COLORS["text"]
                bold = False
                align = PP_ALIGN.LEFT
            for p in cell.text_frame.paragraphs:
                p.alignment = align
                for run in p.runs:
                    set_font(run, font_size, color, bold)


def add_table(slide, x, y, w, h, headers, rows, col_widths, font_size=6):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = shape.table
    for c, header in enumerate(headers):
        table.cell(0, c).text = header
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = width
    style_table(table, font_size=font_size)
    return shape


PROJECT_ROWS = [
    [
        "P01",
        "报表上线与数据确认专项",
        "电流电压、光学报表预计 6/22 上线测试，待确认最终提交数据",
        "完成上线验证，固化最终数据提交口径",
        "测试结论 / 数据清单 / 异常闭环记录",
        "里程碑时间轴",
    ],
    [
        "P02",
        "彩斑/点类不良 AVI 检验与复判标准专项",
        "针对彩斑、点类不良提供 AVI 设备检验原理及算法，及人员复判标准",
        "明确设备检验逻辑与人工复判一致性标准",
        "检验原理 / 算法逻辑 / 复判 SOP",
        "标准化矩阵",
    ],
    [
        "P03",
        "Z571 放量与彩斑拦截优化专项",
        "Z571 项目放量策划，针对上一代彩斑优化拦截措施",
        "支撑 Z571 放量并降低彩斑风险外流",
        "放量计划 / 拦截策略 / 验证结果",
        "阶段甘特图",
    ],
    [
        "P04",
        "CVD Mask 来料刮伤监控与处置专项",
        "提供 CVD Mask 来料针对刮伤监控 rule 和处理方式",
        "形成来料刮伤监控规则与快速处置流程",
        "监控 rule / 判定标准 / 处置流程",
        "流程闭环图",
    ],
    [
        "P05",
        "Particle 改善措施与效果预估专项",
        "罗列 particle 改善措施及对策预估效果，设备优化可单独体现",
        "量化 particle 改善动作及预估收益",
        "措施清单 / 效果预估 / 设备优化项",
        "收益瀑布图",
    ],
]


def build_project_table_slide(slide):
    add_header(
        slide,
        "01 专项项目",
        "图一项次整理为专项项目清单",
        "将原“移班专项 / PT 专项遗留”项次转换为可跟踪、可汇报、可落责的专项项目",
        "01 / 05",
    )
    headers = ["编号", "专项项目", "来源项次", "专项目标", "关键交付物", "建议呈现"]
    rows = PROJECT_ROWS
    add_table(
        slide,
        Inches(0.36),
        Inches(1.48),
        Inches(12.62),
        Inches(4.9),
        headers,
        rows,
        [Inches(0.5), Inches(2.05), Inches(3.55), Inches(2.2), Inches(2.55), Inches(1.15)],
        5,
    )
    add_round_rect(slide, Inches(0.58), Inches(6.52), Inches(11.75), Inches(0.36), COLORS["soft_blue"], COLORS["light_line"])
    add_text(
        slide,
        "填充建议：后续可追加责任部门、计划完成日、当前状态、完成率、风险说明，用于周报/年中总结持续更新。",
        Inches(0.72),
        Inches(6.6),
        Inches(11.4),
        Inches(0.16),
        8,
        COLORS["grey"],
    )


def build_project_board_slide(slide):
    add_header(
        slide,
        "02 专项看板",
        "专项项目推进看板",
        "用卡片形式体现项目类型、关键输出和汇报图表，适合作为项目管理页或汇总页",
        "02 / 05",
    )
    lanes = [
        ("报表/数据", COLORS["soft_blue"], COLORS["blue"], [PROJECT_ROWS[0]]),
        ("检验/标准", COLORS["soft_orange"], COLORS["orange"], [PROJECT_ROWS[1], PROJECT_ROWS[3]]),
        ("放量/改善", COLORS["soft_green"], COLORS["green"], [PROJECT_ROWS[2], PROJECT_ROWS[4]]),
    ]
    x_positions = [Inches(0.62), Inches(4.42), Inches(8.22)]
    for idx, (lane_title, fill, accent, items) in enumerate(lanes):
        x = x_positions[idx]
        add_round_rect(slide, x, Inches(1.48), Inches(3.45), Inches(4.82), fill, COLORS["light_line"])
        add_text(slide, lane_title, x + Inches(0.18), Inches(1.63), Inches(2.8), Inches(0.24), 13, COLORS["text"], True)
        add_rect(slide, x + Inches(0.18), Inches(1.96), Inches(0.55), Inches(0.05), accent, accent)
        for item_idx, row in enumerate(items):
            card_y = Inches(2.18 + item_idx * 1.85)
            add_round_rect(slide, x + Inches(0.2), card_y, Inches(3.05), Inches(1.5), COLORS["white"], COLORS["light_line"])
            add_text(slide, row[0], x + Inches(0.36), card_y + Inches(0.12), Inches(0.45), Inches(0.22), 8, accent, True)
            add_text(slide, row[1], x + Inches(0.86), card_y + Inches(0.1), Inches(2.2), Inches(0.32), 10, COLORS["text"], True)
            add_text(slide, "输出：" + row[4], x + Inches(0.36), card_y + Inches(0.55), Inches(2.75), Inches(0.24), 7, COLORS["grey"])
            add_text(slide, "图表：" + row[5], x + Inches(0.36), card_y + Inches(0.96), Inches(2.75), Inches(0.22), 7, accent, False)
    add_round_rect(slide, Inches(0.68), Inches(6.48), Inches(11.55), Inches(0.38), COLORS["bg"], COLORS["light_line"])
    add_text(slide, "可替换字段：状态（未启动/进行中/已关闭）、责任部门、截止日期、阻塞说明、闭环结论。", Inches(0.84), Inches(6.57), Inches(11.2), Inches(0.17), 8, COLORS["grey"])


def build_cpk_slide(slide, product_name: str, page_no: str, before: float = 1.37, after: float = 1.43):
    add_header(
        slide,
        "03 CPK 优化",
        f"{product_name} CPK 优化幅度",
        "按图二样式复制，产品名称、优化前后数值和提升幅度均可直接替换",
        page_no,
    )
    add_round_rect(slide, Inches(0.58), Inches(1.18), Inches(10.55), Inches(5.55), COLORS["white"], COLORS["light_line"])

    chart_left = Inches(1.82)
    chart_top = Inches(2.08)
    chart_width = Inches(7.75)
    chart_height = Inches(3.75)
    baseline_y = chart_top + chart_height
    y_axis_x = chart_left

    add_line(slide, y_axis_x, chart_top, y_axis_x, baseline_y, COLORS["light_line"], 1.2)
    add_line(slide, y_axis_x, baseline_y, chart_left + chart_width, baseline_y, COLORS["light_line"], 1.2)

    max_val = 1.7
    scale_h = Inches(3.0)
    bar_w = Inches(1.65)
    before_h = scale_h * (before / max_val)
    after_h = scale_h * (after / max_val)
    before_x = Inches(3.25)
    after_x = Inches(6.45)
    before_y = baseline_y - before_h
    after_y = baseline_y - after_h

    add_rect(slide, before_x, before_y, bar_w, before_h, COLORS["grey"], COLORS["grey"])
    add_rect(slide, after_x, after_y, bar_w, after_h, COLORS["green"], COLORS["green"])
    add_text(slide, f"{before:.2f}", before_x, before_y - Inches(0.35), bar_w, Inches(0.22), 11, COLORS["text"], True, PP_ALIGN.CENTER)
    add_text(slide, f"{after:.2f}", after_x, after_y - Inches(0.35), bar_w, Inches(0.22), 11, COLORS["text"], True, PP_ALIGN.CENTER)
    add_text(slide, "优化前", before_x, baseline_y + Inches(0.18), bar_w, Inches(0.24), 10, COLORS["grey"], False, PP_ALIGN.CENTER)
    add_text(slide, "优化后", after_x, baseline_y + Inches(0.18), bar_w, Inches(0.24), 10, COLORS["grey"], False, PP_ALIGN.CENTER)

    lift = (after - before) / before
    add_text(slide, f"提升幅度 {lift:.0%}", Inches(5.1), Inches(1.72), Inches(1.6), Inches(0.26), 11, COLORS["blue"], True, PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(33, Inches(4.95), Inches(2.15), Inches(1.65), Inches(0.48))
    arrow.rotation = -35
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(255, 230, 153)
    arrow.line.color.rgb = COLORS["yellow"]
    arrow.line.width = Pt(1)

    add_round_rect(slide, Inches(9.88), Inches(2.02), Inches(1.75), Inches(1.4), COLORS["soft_blue"], COLORS["light_line"])
    add_text(slide, "产品名称", Inches(10.05), Inches(2.18), Inches(1.4), Inches(0.2), 8, COLORS["grey"], False, PP_ALIGN.CENTER)
    add_text(slide, product_name, Inches(10.05), Inches(2.48), Inches(1.4), Inches(0.32), 14, COLORS["text"], True, PP_ALIGN.CENTER)
    add_text(slide, "备注：可替换为实际产品名", Inches(10.02), Inches(2.94), Inches(1.45), Inches(0.2), 7, COLORS["grey"], False, PP_ALIGN.CENTER)

    add_round_rect(slide, Inches(9.88), Inches(3.76), Inches(1.75), Inches(1.52), COLORS["soft_green"], COLORS["light_line"])
    add_text(slide, "数据口径", Inches(10.05), Inches(3.9), Inches(1.4), Inches(0.2), 8, COLORS["grey"], False, PP_ALIGN.CENTER)
    add_text(slide, "CPK", Inches(10.05), Inches(4.18), Inches(1.4), Inches(0.26), 13, COLORS["green"], True, PP_ALIGN.CENTER)
    add_text(slide, "优化前/优化后", Inches(10.02), Inches(4.64), Inches(1.45), Inches(0.2), 7, COLORS["grey"], False, PP_ALIGN.CENTER)

    add_round_rect(slide, Inches(0.84), Inches(6.28), Inches(11.0), Inches(0.34), COLORS["bg"], COLORS["light_line"])
    add_text(slide, "填充提示：将 1.37、1.43 和提升幅度替换为该产品真实 CPK 数据；需要更多产品时可复制本页。", Inches(1.0), Inches(6.36), Inches(10.6), Inches(0.16), 8, COLORS["grey"])


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    build_project_table_slide(slide)

    slide = prs.slides.add_slide(blank)
    build_project_board_slide(slide)

    for page_no, product in zip(["03 / 05", "04 / 05", "05 / 05"], ["产品 A", "产品 B", "产品 C"]):
        slide = prs.slides.add_slide(blank)
        build_cpk_slide(slide, product, page_no)

    output = EXPORT_DIR / f"special_projects_cpk_products_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build_deck())
