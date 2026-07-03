from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_DIR = Path(r"D:\wzy\Python\vivo-project\projects\customer_midyear_table_templates")
EXPORT_DIR = PROJECT_DIR / "exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


COLORS = {
    "bg": RGBColor(247, 250, 253),
    "white": RGBColor(255, 255, 255),
    "blue": RGBColor(91, 155, 213),
    "deep_blue": RGBColor(68, 114, 196),
    "navy": RGBColor(31, 78, 121),
    "orange": RGBColor(237, 125, 49),
    "green": RGBColor(112, 173, 71),
    "yellow": RGBColor(255, 192, 0),
    "grey": RGBColor(112, 120, 132),
    "light_grey": RGBColor(231, 236, 244),
    "mid_grey": RGBColor(205, 216, 230),
    "text": RGBColor(68, 84, 106),
}


def set_font(run, size: int, color=COLORS["text"], bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def set_text(shape, text: str, size: int = 12, color=COLORS["text"], bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)


def add_text(slide, text: str, x, y, w, h, size=12, color=COLORS["text"], bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, w, h)
    set_text(shape, text, size, color, bold, align)
    return shape


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.6)
    return shape


def add_header(slide, section: str, title: str, subtitle: str, page_no: str) -> None:
    add_rect(slide, Inches(0), Inches(0), Inches(0.16), Inches(7.5), COLORS["blue"], COLORS["blue"])
    add_rect(slide, Inches(0.16), Inches(0.88), Inches(0.06), Inches(1.3), COLORS["orange"], COLORS["orange"])
    add_rect(slide, Inches(0.16), Inches(2.34), Inches(0.06), Inches(1.3), COLORS["green"], COLORS["green"])
    add_text(slide, section, Inches(0.58), Inches(0.17), Inches(2.2), Inches(0.25), 9, COLORS["blue"], True)
    add_text(slide, title, Inches(0.58), Inches(0.42), Inches(7.8), Inches(0.42), 21, COLORS["text"], True)
    add_text(slide, subtitle, Inches(0.58), Inches(0.83), Inches(8.4), Inches(0.24), 8, COLORS["grey"])
    line = slide.shapes.add_connector(1, Inches(0.58), Inches(1.12), Inches(12.4), Inches(1.12))
    line.line.color.rgb = COLORS["mid_grey"]
    line.line.width = Pt(1)
    add_text(slide, page_no, Inches(12.25), Inches(6.95), Inches(0.55), Inches(0.22), 8, COLORS["grey"], False, PP_ALIGN.RIGHT)


def style_table(table, header_fill=COLORS["deep_blue"], header_text=COLORS["white"], font_size=8) -> None:
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_fill
                color = header_text
                bold = True
            else:
                cell.fill.fore_color.rgb = COLORS["white"] if row_idx % 2 else COLORS["bg"]
                color = COLORS["text"]
                bold = False
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if row_idx == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_font(run, font_size, color, bold)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=8, header_fill=COLORS["deep_blue"]):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = shape.table
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = value
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = width
    style_table(table, header_fill=header_fill, font_size=font_size)
    return shape


def add_tip(slide, text: str) -> None:
    add_rect(slide, Inches(0.58), Inches(6.5), Inches(11.75), Inches(0.38), COLORS["light_grey"], COLORS["mid_grey"])
    add_text(slide, text, Inches(0.72), Inches(6.58), Inches(11.45), Inches(0.17), 8, COLORS["grey"])


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_header(slide, "00 表格索引", "客户专项年中总结：建议数据表模板", "独立 PPT，可直接复制表格或填入真实数据后支撑图表页", "01 / 06")
    headers = ["模块", "建议表格", "支撑图表/页面", "关键字段", "口径备注"]
    rows = [
        ["系统搭建", "系统搭建里程碑表", "进度时间轴", "阶段、任务、完成率、责任人", "按周/月更新"],
        ["系统搭建", "系统模块建设清单", "模块状态卡", "模块、状态、接口/规则数", "状态口径统一"],
        ["风险品处置", "风险品处置闭环明细表", "处置明细/责任闭环", "状态、责任部门、截止日期、闭环说明", "对应截图提示"],
        ["风险品处置", "风险品处置统计汇总表", "KPI 卡片/漏斗图", "预警数、风险品数、处置数、闭环率", "用于领导页"],
        ["制程提升", "制程提升改善项目表", "CPK 对比/进展矩阵", "优化前后、提升幅度、固化输出", "体现改善收益"],
    ]
    add_table(slide, Inches(0.58), Inches(1.42), Inches(12.1), Inches(4.15), headers, rows, [Inches(1.15), Inches(2.15), Inches(2.0), Inches(3.6), Inches(2.9)], 8)
    add_tip(slide, "建议：先用这些表统一数据口径，再回填到三页汇报模板中的图表占位。")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "01 系统搭建", "系统搭建里程碑表", "用于填充系统搭建进展时间轴，呈现阶段、进度、责任和下一步", "02 / 06")
    headers = ["阶段", "关键任务", "计划完成", "实际进度", "完成率", "责任人", "风险/阻塞", "下一步"]
    rows = [
        ["需求梳理", "业务口径与数据范围确认", "YYYY-MM-DD", "已完成/进行中", "XX%", "XX", "待填充", "待填充"],
        ["规则配置", "预警规则、阈值、产品清单配置", "YYYY-MM-DD", "已完成/进行中", "XX%", "XX", "待填充", "待填充"],
        ["数据接入", "数据源、接口、权限联通", "YYYY-MM-DD", "已完成/进行中", "XX%", "XX", "待填充", "待填充"],
        ["看板联调", "报表看板、筛选条件、权限校验", "YYYY-MM-DD", "已完成/进行中", "XX%", "XX", "待填充", "待填充"],
        ["试运行迭代", "用户反馈、问题修复、上线评审", "YYYY-MM-DD", "计划中/进行中", "XX%", "XX", "待填充", "待填充"],
    ]
    add_table(slide, Inches(0.48), Inches(1.42), Inches(12.35), Inches(4.75), headers, rows, [Inches(1.0), Inches(2.55), Inches(1.25), Inches(1.35), Inches(0.78), Inches(0.78), Inches(2.15), Inches(2.0)], 7)
    add_tip(slide, "图表映射：阶段 + 计划完成 + 实际进度可生成里程碑时间轴；完成率可生成进度条。")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "01 系统搭建", "系统模块建设清单", "用于填充关键模块状态卡，展示系统能力沉淀和建设覆盖范围", "03 / 06")
    headers = ["模块名称", "建设内容", "当前状态", "接口/规则数", "上线范围", "负责人", "备注"]
    rows = [
        ["数据源接入", "接入 XX 类数据源，覆盖 XX 条记录", "待填充", "XX", "XX 产品/线体", "XX", "待填充"],
        ["预警规则配置", "配置风险触发规则、阈值和分层口径", "待填充", "XX", "XX 产品/线体", "XX", "待填充"],
        ["处置流转闭环", "建立责任分派、跟进、闭环说明字段", "待填充", "XX", "XX 部门", "XX", "待填充"],
        ["报表看板发布", "输出日报/周报/月报与专项跟踪看板", "待填充", "XX", "XX 用户", "XX", "待填充"],
        ["权限与运维", "配置账号权限、操作记录、问题响应机制", "待填充", "XX", "XX 角色", "XX", "待填充"],
    ]
    add_table(slide, Inches(0.58), Inches(1.42), Inches(12.1), Inches(4.75), headers, rows, [Inches(1.35), Inches(3.2), Inches(1.1), Inches(1.05), Inches(1.7), Inches(0.9), Inches(2.8)], 7)
    add_tip(slide, "图表映射：当前状态可生成模块状态卡；接口/规则数可生成小型 KPI。")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "02 风险品处置", "风险品处置闭环明细表", "对应“已关闭 / 持续跟进 / 升级专项 / 责任部门 / 截止日期 / 闭环说明”", "04 / 06")
    headers = ["序号", "产品/型号", "预警来源", "预警数", "风险等级", "处置状态", "责任部门", "截止日期", "闭环说明"]
    rows = [
        ["1", "产品 A", "规则/人工/系统", "XX", "高/中/低", "已关闭", "XX 部门", "YYYY-MM-DD", "已完成 XX 措施，结果待填充"],
        ["2", "产品 B", "规则/人工/系统", "XX", "高/中/低", "持续跟进", "XX 部门", "YYYY-MM-DD", "当前阻塞/下一动作待填充"],
        ["3", "产品 C", "规则/人工/系统", "XX", "高/中/低", "升级专项", "XX 部门", "YYYY-MM-DD", "升级原因与专项编号待填充"],
        ["4", "产品 D", "规则/人工/系统", "XX", "高/中/低", "待判定", "XX 部门", "YYYY-MM-DD", "待填充"],
        ["5", "产品 E", "规则/人工/系统", "XX", "高/中/低", "待判定", "XX 部门", "YYYY-MM-DD", "待填充"],
    ]
    add_table(slide, Inches(0.36), Inches(1.42), Inches(12.62), Inches(4.75), headers, rows, [Inches(0.45), Inches(1.25), Inches(1.35), Inches(0.72), Inches(0.85), Inches(1.08), Inches(1.05), Inches(1.2), Inches(4.67)], 6)
    add_tip(slide, "建议状态枚举：已关闭、持续跟进、升级专项、待判定。这样可直接统计漏斗和闭环率。")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "02 风险品处置", "风险品处置统计汇总表", "用于填充 KPI 卡片、漏斗图和闭环率说明", "05 / 06")
    headers = ["统计项", "数量", "占比/转化率", "环比/同比", "口径说明", "图表用途"]
    rows = [
        ["触发预警数量", "XX", "-", "XX%", "统计期内触发的全部预警条数", "顶部 KPI / 漏斗第一层"],
        ["风险产品数", "XX", "XX%", "XX%", "去重后的风险产品/型号数量", "顶部 KPI / 漏斗第二层"],
        ["已关闭", "XX", "XX%", "XX%", "已完成处置并有闭环说明", "处置分类"],
        ["持续跟进", "XX", "XX%", "XX%", "仍在责任部门推进", "处置分类"],
        ["升级专项", "XX", "XX%", "XX%", "需专项机制跟踪的风险品", "处置分类"],
        ["产品处置数", "XX", "XX%", "XX%", "已关闭 + 已进入专项闭环口径", "顶部 KPI / 漏斗第三层"],
        ["闭环率", "XX%", "XX%", "XX%", "产品处置数 / 风险产品数", "漏斗转化率"],
    ]
    add_table(slide, Inches(0.54), Inches(1.34), Inches(12.15), Inches(5.1), headers, rows, [Inches(1.55), Inches(0.85), Inches(1.1), Inches(1.0), Inches(4.35), Inches(3.3)], 7, COLORS["orange"])
    add_tip(slide, "图表映射：数量字段生成 KPI；统计项顺序可直接作为漏斗层级和处置分类。")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 制程提升", "制程提升改善项目表", "用于填充 CPK 优化、流程优化和专项进展矩阵", "06 / 06")
    headers = ["序号", "专项/工序", "问题点", "优化动作", "CPK 优化前", "CPK 优化后", "提升幅度", "当前进度", "固化输出"]
    rows = [
        ["1", "专项 A / 工序 A", "待填充", "待填充", "XX", "XX", "XX%", "XX%", "SOP/规则/看板"],
        ["2", "专项 B / 工序 B", "待填充", "待填充", "XX", "XX", "XX%", "XX%", "SOP/规则/看板"],
        ["3", "专项 C / 工序 C", "待填充", "待填充", "XX", "XX", "XX%", "XX%", "SOP/规则/看板"],
        ["4", "流程优化项 D", "待填充", "待填充", "-", "-", "XX%", "XX%", "流程文件/培训记录"],
        ["5", "流程优化项 E", "待填充", "待填充", "-", "-", "XX%", "XX%", "流程文件/培训记录"],
    ]
    add_table(slide, Inches(0.34), Inches(1.42), Inches(12.65), Inches(4.75), headers, rows, [Inches(0.42), Inches(1.55), Inches(1.7), Inches(2.25), Inches(1.0), Inches(1.0), Inches(0.9), Inches(0.9), Inches(2.93)], 6, COLORS["green"])
    add_tip(slide, "图表映射：CPK 优化前后生成对比柱状图；当前进度生成专项进度条；固化输出体现改善沉淀。")

    output = EXPORT_DIR / f"customer_midyear_table_templates_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build_deck())
